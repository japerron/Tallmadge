"""
renderers/pptx.py — PowerPoint presentation generator
Creates a branded .pptx with title slide + one slide per prompt.
Uses python-pptx.
"""

import io
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json

from core.results import parse_grid_data, format_sql, strip_markdown

OUTPUT_BASE = Path(__file__).parent.parent / "output"

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)


def _rgb(hex_color: str) -> RGBColor:
    c = hex_color.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _add_rect(slide, left, top, width, height, fill_hex, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_text_box(slide, text, left, top, width, height,
                  font_name="Calibri", size=12, bold=False, italic=False,
                  color="#000000", align=PP_ALIGN.LEFT, wrap=True,
                  line_spacing=None) -> None:
    from pptx.util import Pt as _Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    para  = tf.paragraphs[0]
    para.alignment = align
    if line_spacing is not None:
        para.line_spacing = line_spacing   # float multiplier e.g. 1.5
    run   = para.add_run()
    run.text = str(text) if text else ""
    run.font.name   = font_name
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


def _style_cell(cell, text, font_name, size, bold=False,
                color="#000000", bg_hex=None):
    """Set text + formatting on a pptx table cell."""
    cell.text = str(text) if text is not None else ""
    if bg_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(bg_hex)
    tf = cell.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name  = font_name
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = _rgb(color)


def _truncate(text: str | None, max_chars: int = 400) -> str:
    if not text:
        return ""
    text = str(text)
    if len(text) > max_chars:
        return text[:max_chars].rsplit(" ", 1)[0] + "…"
    return text


def _clean_content(text: str | None) -> str:
    """
    Collapse blank lines and normalise whitespace for slide display.
    Joins non-empty lines with a single space so the text box wraps cleanly
    without wasting lines on paragraph gaps.
    """
    if not text:
        return ""
    return " ".join(line.strip() for line in str(text).splitlines() if line.strip())


def _truncate_lines(text: str | None, box_w_in: float, pt: float,
                    max_lines: int) -> str:
    """
    Truncate text to at most max_lines visual lines in a box of box_w_in inches
    at font size pt, using the same proportional-font heuristic as _fit_pt.
    Appends '…' when text is cut short.
    """
    if not text:
        return ""
    text = str(text)
    chars_per_line = max(1.0, box_w_in * _CHARS_PER_INCH_10PT * (10.0 / pt))

    accepted: list[str] = []
    line_len:  float    = 0.0
    lines_used: int     = 1

    for i, word in enumerate(text.split()):
        # cost of adding this word to the current line (leading space if not first)
        cost = len(word) + (1 if accepted and line_len > 0 else 0)
        if accepted and line_len + cost > chars_per_line:
            lines_used += 1
            if lines_used > max_lines:
                return " ".join(accepted) + "…"
            line_len = len(word)
        else:
            line_len += cost
        accepted.append(word)

    return " ".join(accepted)


# ── Auto font-size helpers ──────────────────────────────────────────────────

_CHARS_PER_INCH_10PT = 13.5   # Calibri/Arial proportional-font estimate at 10 pt
_LINE_SPACING        = 1.30   # line-height multiplier

def _fit_pt(text: str | None, box_w_in: float, box_h_in: float,
            max_pt: float = 13.0, min_pt: float = 8.5) -> float:
    """
    Return the largest font size in [min_pt, max_pt] (half-point steps) that
    approximately fits `text` into a box of box_w_in × box_h_in inches.
    Uses a proportional-font character-count heuristic.
    """
    if not text:
        return max_pt
    n  = len(text)
    pt = max_pt
    while pt >= min_pt:
        chars_per_line = box_w_in * _CHARS_PER_INCH_10PT * (10.0 / pt)
        lines_avail    = box_h_in / ((pt / 72.0) * _LINE_SPACING)
        if (n / max(chars_per_line, 1)) <= lines_avail:
            return round(pt, 1)
        pt -= 0.5
    return min_pt


def _grid_font(n_cols: int, is_header: bool) -> float:
    """Scale table font size with column count: fewer cols → bigger text."""
    base = 10.0 if is_header else 9.5
    # Each extra column beyond 4 reduces by 0.5 pt; floor at 7.5 pt
    return max(7.5, base - max(0, n_cols - 4) * 0.5)


def render(envelope: dict, style: dict, out_dir: Path | None = None) -> Path:
    """
    Generate a PowerPoint presentation from results envelope.
    Returns the path to the saved .pptx file.
    """
    results  = envelope.get("results", [])
    meta     = envelope.get("meta", {})
    _raw_date = meta.get("runDate", "")
    run_date  = _raw_date[:10]                                        # YYYY-MM-DD  (display)
    run_ts    = (_raw_date[:19].replace("T", "_").replace(":", "-")   # YYYY-MM-DD_HH-MM-SS
                 if len(_raw_date) >= 19 else run_date)               # (filename)
    mode     = meta.get("mode", "standard").title()

    if out_dir is None:
        out_dir = OUTPUT_BASE
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"tallmadge_{run_ts}.pptx"

    primary   = style["primary"]
    secondary = style["secondary"]
    accent    = style["accent"]
    font      = style["font"]

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── Title slide ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)

    # Full background
    _add_rect(slide, 0, 0, W, H, primary)

    # Accent bar (left 40%)
    _add_rect(slide, 0, 0, Inches(5.3), H, secondary)

    # Title
    _add_text_box(slide, "Agent Test Results",
                  Inches(5.6), Inches(2.2), Inches(7), Inches(1.2),
                  font_name=font, size=36, bold=True, color="#FFFFFF",
                  align=PP_ALIGN.LEFT)

    # Subtitle
    subtitle = f"{style['name']}  ·  {run_date}  ·  {mode} Mode"
    _add_text_box(slide, subtitle,
                  Inches(5.6), Inches(3.5), Inches(7), Inches(0.6),
                  font_name=font, size=14, color="#FFFFFF", align=PP_ALIGN.LEFT)

    # Stats bar at bottom
    stats = f"✓ {meta.get('successful',0)} Success   ✗ {meta.get('errors',0)} Errors   {meta.get('totalPrompts',0)} Total Prompts"
    _add_text_box(slide, stats,
                  Inches(5.6), Inches(6.5), Inches(7), Inches(0.5),
                  font_name=font, size=11, color="#FFFFFF", italic=True, align=PP_ALIGN.LEFT)

    # Brand name on the left panel
    _add_text_box(slide, "🔮",
                  Inches(1.5), Inches(2.8), Inches(2.5), Inches(1),
                  font_name=font, size=48, align=PP_ALIGN.CENTER)
    _add_text_box(slide, style["name"],
                  Inches(0.5), Inches(4.0), Inches(4.3), Inches(0.6),
                  font_name=font, size=13, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

    # ── One slide per prompt ───────────────────────────────────────────────────
    for r in results:
        slide = prs.slides.add_slide(blank_layout)

        # Top header bar
        _add_rect(slide, 0, 0, W, Inches(1.1), primary)

        # Prompt number
        _add_text_box(slide, str(r["id"]),
                      Inches(0.3), Inches(0.15), Inches(0.8), Inches(0.8),
                      font_name=font, size=22, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)

        # Category label
        _add_text_box(slide, r["category"].upper(),
                      Inches(1.2), Inches(0.15), Inches(11), Inches(0.35),
                      font_name=font, size=9, color="#FFFFFF", italic=True)

        # Prompt text
        _add_text_box(slide, r["prompt"],
                      Inches(1.2), Inches(0.45), Inches(11.5), Inches(0.55),
                      font_name=font, size=15, bold=True, color="#FFFFFF")

        # Status pill
        status_color = "#1B6B2F" if r["status"] == "Success" else "#B71C1C"
        status_bg    = "#E8F5E9" if r["status"] == "Success" else "#FDECEA"
        _add_rect(slide, Inches(11.3), Inches(0.2), Inches(1.7), Inches(0.5), status_bg)
        _add_text_box(slide, f"{'✓' if r['status']=='Success' else '✗'} {r['status']}",
                      Inches(11.3), Inches(0.2), Inches(1.7), Inches(0.5),
                      font_name=font, size=9, bold=True, color=status_color, align=PP_ALIGN.CENTER)

        # ── Content area: two columns ──────────────────────────────────────────
        # Left column: Response Text + Interpretation
        y = Inches(1.2)
        left_x = Inches(0.3)
        col_w  = Inches(6.2)

        resp_text = _truncate_lines(_clean_content(r.get("responseText")),   6.2, 12, 6)
        _add_text_box(slide, "RESPONSE",
                      left_x, y, col_w, Inches(0.25),
                      font_name=font, size=10, bold=True, color=primary)
        y += Inches(0.28)
        _add_text_box(slide, resp_text,
                      left_x, y, col_w, Inches(1.8),
                      font_name=font, size=12, italic=True, color="#555555",
                      wrap=True, line_spacing=1.5)
        y += Inches(1.9)

        if r.get("interpretedQuestion"):
            interp_text = _truncate_lines(_clean_content(r.get("interpretedQuestion")), 6.2, 12, 3)
            _add_text_box(slide, "INTERPRETATION",
                          left_x, y, col_w, Inches(0.25),
                          font_name=font, size=10, bold=True, color=primary)
            y += Inches(0.28)
            _add_text_box(slide, interp_text,
                          left_x, y, col_w, Inches(1.0),
                          font_name=font, size=12, italic=True, color="#555555",
                          wrap=True, line_spacing=1.5)
            y += Inches(1.1)

        if r.get("explanation"):
            expl_text = _truncate_lines(_clean_content(strip_markdown(r.get("explanation"))), 6.2, 12, 4)
            _add_text_box(slide, "EXPLANATION",
                          left_x, y, col_w, Inches(0.25),
                          font_name=font, size=10, bold=True, color=primary)
            y += Inches(0.28)
            _add_text_box(slide, expl_text,
                          left_x, y, col_w, Inches(1.2),
                          font_name=font, size=12, italic=True, color="#555555",
                          wrap=True, line_spacing=1.5)

        # Right column: image (2/3 height if present) + insights + SQL
        # SQL shown only when image OR insights is absent; never truncated.
        # WHERE tokens never shown.
        ry             = Inches(1.2)
        right_x        = Inches(6.8)
        right_w        = Inches(6.2)
        col_bottom     = Inches(6.9)          # footer top edge
        available_h    = col_bottom - ry       # Inches(5.7)

        has_image    = bool(r.get("imagesData"))
        has_insights = bool(r.get("insights"))
        has_sql      = bool(r.get("sqlQueries"))
        show_sql     = has_sql and not (has_image and has_insights)

        # ── Image (2/3 of available height) ───────────────────────────────────
        if has_image:
            img_h = available_h * 2 / 3
            img_w = img_h * 800 / 600
            if img_w > right_w:          # constrain width if aspect pushes it wide
                img_w = right_w
                img_h = img_w * 600 / 800
            try:
                img        = r["imagesData"][0]
                img_stream = io.BytesIO(base64.b64decode(img["data"]))
                slide.shapes.add_picture(
                    img_stream,
                    right_x + (right_w - img_w) / 2,
                    ry,
                    img_w, img_h,
                )
            except Exception:
                pass
            ry += img_h + Inches(0.1)

        # ── Insights (shown whenever present) ─────────────────────────────────
        if has_insights:
            remaining    = col_bottom - ry
            # When SQL also follows, reserve space for it; otherwise take all
            insights_h   = (Inches(1.6) if show_sql
                            else remaining - Inches(0.3))
            insights_h   = max(insights_h, Inches(0.5))
            _add_text_box(slide, "INSIGHTS",
                          right_x, ry, right_w, Inches(0.25),
                          font_name=font, size=10, bold=True, color=primary)
            ry += Inches(0.28)
            _add_text_box(slide, _truncate_lines(_clean_content(r.get("insights")), 6.2, 12, 5),
                          right_x, ry, right_w, insights_h,
                          font_name=font, size=12, italic=True,
                          color="#555555", wrap=True, line_spacing=1.5)
            ry += insights_h + Inches(0.1)

        # ── SQL (shown when image or insights absent; never truncated) ─────────
        if show_sql:
            sql_raw  = (r["sqlQueries"][0] if isinstance(r["sqlQueries"], list)
                        else r["sqlQueries"])
            sql_text = _clean_content(format_sql(str(sql_raw)))
            remaining = col_bottom - ry
            sql_h     = max(remaining - Inches(0.3), Inches(0.5))
            _add_text_box(slide, "SQL",
                          right_x, ry, right_w, Inches(0.25),
                          font_name=font, size=10, bold=True, color=primary)
            ry += Inches(0.28)
            _add_text_box(slide, sql_text,
                          right_x, ry, right_w, sql_h,
                          font_name=font, size=10, italic=True,
                          color="#555555", wrap=True, line_spacing=1.5)

        # Attributes / Metrics footer row
        attrs   = ", ".join(r.get("attributesUsed") or []) or "—"
        metrics = ", ".join(r.get("metricsUsed") or [])     or "—"
        rt_str  = f"{r.get('responseTime', '—')}s"

        _add_rect(slide, 0, Inches(6.9), W, Inches(0.6), "EEEEEE")
        _add_text_box(slide, f"🏷 {_truncate(attrs, 80)}",
                      Inches(0.3), Inches(6.92), Inches(5.5), Inches(0.5),
                      font_name=font, size=8, color="#555555")
        _add_text_box(slide, f"📐 {_truncate(metrics, 80)}",
                      Inches(5.8), Inches(6.92), Inches(5.5), Inches(0.5),
                      font_name=font, size=8, color="#555555")
        _add_text_box(slide, f"⏱ {rt_str}",
                      Inches(11.8), Inches(6.92), Inches(1.2), Inches(0.5),
                      font_name=font, size=8, color="#888888", align=PP_ALIGN.RIGHT)

        # ── Optional data slide for gridData ───────────────────────────────────
        g_headers, g_rows = parse_grid_data(r.get("gridData"))
        if g_headers and g_rows:
            MAX_COLS   = 8
            MAX_ROWS   = 18
            total_cols = len(g_headers)
            total_rows = len(g_rows)
            g_headers  = g_headers[:MAX_COLS]
            g_rows     = g_rows[:MAX_ROWS]
            n_cols     = len(g_headers)
            n_rows     = len(g_rows)

            dslide = prs.slides.add_slide(blank_layout)

            # Header bar (same style as content slides)
            _add_rect(dslide, 0, 0, W, Inches(1.1), primary)
            _add_text_box(dslide, str(r["id"]),
                          Inches(0.3), Inches(0.15), Inches(0.8), Inches(0.8),
                          font_name=font, size=22, bold=True, color="#FFFFFF",
                          align=PP_ALIGN.CENTER)
            _add_text_box(dslide, f"{r['category'].upper()} — DATA",
                          Inches(1.2), Inches(0.15), Inches(11), Inches(0.35),
                          font_name=font, size=9, color="#FFFFFF", italic=True)
            _add_text_box(dslide, r["prompt"],
                          Inches(1.2), Inches(0.45), Inches(11.5), Inches(0.55),
                          font_name=font, size=15, bold=True, color="#FFFFFF")

            # Grid table
            tbl_left = Inches(0.3)
            tbl_top  = Inches(1.3)
            tbl_w    = W - Inches(0.6)
            tbl_h    = Inches(5.5)
            col_w    = int(tbl_w / n_cols)

            table = dslide.shapes.add_table(
                n_rows + 1, n_cols, tbl_left, tbl_top, tbl_w, tbl_h
            ).table

            # Uniform column widths
            for ci in range(n_cols):
                table.columns[ci].width = col_w

            # Header row
            hdr_pt = _grid_font(n_cols, is_header=True)
            dat_pt = _grid_font(n_cols, is_header=False)
            table.rows[0].height = Inches(0.32)
            for ci, h in enumerate(g_headers):
                _style_cell(table.cell(0, ci), h, font, hdr_pt,
                            bold=True, color="#FFFFFF", bg_hex=primary)

            # Data rows
            row_h = int(tbl_h / (n_rows + 1))
            for ri, row in enumerate(g_rows):
                table.rows[ri + 1].height = row_h
                bg = "F5F5F5" if ri % 2 == 0 else "FFFFFF"
                vals = row if isinstance(row, (list, tuple)) else list(row.values())
                for ci, val in enumerate(vals[:n_cols]):
                    _style_cell(table.cell(ri + 1, ci), val, font, dat_pt,
                                color="#333333", bg_hex=bg)

            # Footer
            _add_rect(dslide, 0, Inches(6.9), W, Inches(0.6), "EEEEEE")
            rows_note = f"{n_rows} of {total_rows} rows" if total_rows > MAX_ROWS else f"{n_rows} rows"
            cols_note = f"{n_cols} of {total_cols} cols" if total_cols > MAX_COLS else f"{n_cols} cols"
            trunc     = " (truncated)" if total_rows > MAX_ROWS or total_cols > MAX_COLS else ""
            _add_text_box(dslide,
                          f"Data — {rows_note} × {cols_note}{trunc}",
                          Inches(0.3), Inches(6.92), Inches(12), Inches(0.5),
                          font_name=font, size=8, color="#555555")

    prs.save(str(out_path))
    return out_path
